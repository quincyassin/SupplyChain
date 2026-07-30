#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成《分单宝-产品功能介绍》PPT。@author huangxinsong"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


ROOT = Path(__file__).resolve().parents[1]
SHOT_DIR = ROOT / "docs" / "screenshots"
OUT_PATH = ROOT / "docs" / "分单宝-产品功能介绍.pptx"

# 配色
C_BG = RGBColor(0xF7, 0xF8, 0xFA)
C_TITLE = RGBColor(0x1A, 0x1A, 0x1A)
C_SUB = RGBColor(0x55, 0x55, 0x55)
C_ACCENT = RGBColor(0x1F, 0x6F, 0x8B)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD = RGBColor(0xE8, 0xF1, 0xF5)
C_LINE = RGBColor(0xD0, 0xD7, 0xDE)


def set_run(run, size=18, bold=False, color=C_TITLE, font="微软雅黑"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=C_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=C_TITLE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"• {item}"
        set_run(run, size=size, color=C_SUB)
    return box


def add_image_fit(slide, image_path: Path, left, top, max_w, max_h):
    if not image_path.exists():
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_w, max_h
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = C_CARD
        placeholder.line.color.rgb = C_LINE
        tf = placeholder.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"截图缺失\n{image_path.name}"
        set_run(run, size=14, color=C_SUB)
        return
    pic = slide.shapes.add_picture(str(image_path), left, top)
    # 等比缩放到框内
    avail_w = max_w
    avail_h = max_h
    ratio = min(avail_w / pic.width, avail_h / pic.height)
    pic.width = int(pic.width * ratio)
    pic.height = int(pic.height * ratio)
    pic.left = int(left + (avail_w - pic.width) / 2)
    pic.top = int(top + (avail_h - pic.height) / 2)


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x1F, 0x6F, 0x8B))
    add_textbox(
        slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
        "分单宝", size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
        "产品功能介绍", size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
        "多平台订单导入 · 智能分单 · 批量回单 · 对账售后一站完成",
        size=18, color=RGBColor(0xD6, 0xEB, 0xF2), align=PP_ALIGN.CENTER,
    )


def overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.6),
                "一、产品概述", size=28, bold=True)
    add_bullets(
        slide, Inches(0.6), Inches(1.2), Inches(12), Inches(5.5),
        [
            "面向电商运营、代发与分销场景的订单处理工具。",
            "支持淘宝、拼多多等多平台 Excel 订单统一导入。",
            "按商家规则自动分单，完成物流回单、对账导出与售后跟踪。",
            "本地 Web 应用：Mac / Windows 单机版双击即可使用。",
            "数据保存在本机，无需额外安装 Java 或数据库环境。",
        ],
        size=20,
    )


def value_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.6),
                "二、核心价值", size=28, bold=True)
    values = [
        ("多平台兼容", "预配置表头映射 + 智能列识别，上传即可归类入库"),
        ("智能分单", "按商家关键字自动分配，未识别进「未定义」可重分"),
        ("全流程闭环", "导入 → 分单 → 回单 → 对账 → 售后一站完成"),
        ("经营数据可见", "维护成本/供货价后，首页按日汇总营业额与利润"),
        ("售后与数据治理", "售后台账、回收站恢复、历史归档清理"),
        ("零门槛部署", "单机版内置运行环境，数据目录可备份迁移"),
    ]
    for i, (title, desc) in enumerate(values):
        row, col = divmod(i, 2)
        left = Inches(0.5 + col * 6.3)
        top = Inches(1.2 + row * 1.8)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.0), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_LINE
        add_textbox(slide, left + Inches(0.25), top + Inches(0.25), Inches(5.5), Inches(0.4),
                    f"{i + 1}. {title}", size=18, bold=True, color=C_ACCENT)
        add_textbox(slide, left + Inches(0.25), top + Inches(0.75), Inches(5.5), Inches(0.55),
                    desc, size=14, color=C_SUB)


def panorama_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.6),
                "三、功能全景", size=28, bold=True)
    modules = [
        "首页（订单处理）", "对账", "售后", "回收站",
        "商品价格维护", "系统配置", "使用手册",
    ]
    for i, name in enumerate(modules):
        left = Inches(0.5 + (i % 4) * 3.15)
        top = Inches(1.5 + (i // 4) * 2.2)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.95), Inches(1.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = C_ACCENT if i < 4 else RGBColor(0x2C, 0x3E, 0x50)
        card.line.fill.background()
        add_textbox(
            slide, left, top + Inches(0.55), Inches(2.95), Inches(0.7),
            name, size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
        )


def step_slide(prs, step_no: int, title: str, bullets: list[str], image_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_textbox(
        slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
        f"步骤 {step_no}｜{title}", size=24, bold=True,
    )
    add_bullets(slide, Inches(0.5), Inches(0.95), Inches(4.6), Inches(5.8), bullets, size=15)
    add_image_fit(
        slide,
        SHOT_DIR / image_name,
        Inches(5.3), Inches(0.95), Inches(7.5), Inches(5.9),
    )


def flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_textbox(slide, Inches(0.5), Inches(0.35), Inches(12), Inches(0.6),
                "典型业务流程", size=28, bold=True)
    steps = [
        "初次配置", "导入订单", "智能分单", "回单导出",
        "价格利润", "对账结算", "售后处理", "数据治理",
    ]
    for i, name in enumerate(steps):
        left = Inches(0.4 + (i % 4) * 3.2)
        top = Inches(1.6 + (i // 4) * 2.5)
        oval = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.9), Inches(1.5)
        )
        oval.fill.solid()
        oval.fill.fore_color.rgb = C_WHITE
        oval.line.color.rgb = C_ACCENT
        add_textbox(
            slide, left, top + Inches(0.25), Inches(2.9), Inches(0.4),
            f"0{i + 1}", size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left, top + Inches(0.7), Inches(2.9), Inches(0.5),
            name, size=18, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER,
        )


def end_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x1F, 0x6F, 0x8B))
    add_textbox(
        slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
        "感谢观看", size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
        "分单宝 · 让订单处理更高效", size=18, color=RGBColor(0xD6, 0xEB, 0xF2), align=PP_ALIGN.CENTER,
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover_slide(prs)
    overview_slide(prs)
    value_slide(prs)
    panorama_slide(prs)
    flow_slide(prs)

    steps = [
        (1, "初次配置：表头映射", [
            "进入「系统配置 → 表头映射」。",
            "配置各平台 Excel 列与系统字段对应关系。",
            "上传订单时依赖此配置识别平台并入库。",
        ], "06-表头映射.png"),
        (2, "商家配置", [
            "维护商家名称与商品关键字。",
            "分单时按关键字匹配商家。",
            "识别词越完整，「未定义」订单越少。",
        ], "08-商家配置.png"),
        (3, "字段映射 / 导出配置", [
            "字段映射：自定义界面字段显示名称。",
            "导出配置：本地目录或浏览器下载。",
            "单机版支持本机选择文件夹。",
        ], "09-导出配置.png"),
        (4, "导入订单（首页）", [
            "点击「上传订单 Excel」。",
            "系统按表头映射自动识别平台并入库。",
            "无法匹配商家规则的订单进入「未定义」。",
        ], "01-首页.png"),
        (5, "按商家分单", [
            "确认分单日期范围后执行「按商家分单」。",
            "已正确分单的订单不会被改动。",
            "按商家生成 Excel，保存位置由导出配置决定。",
        ], "13-首页操作台.png"),
        (6, "填写物流与回单导出", [
            "批量粘贴系统编号、快递公司、快递单号。",
            "订单标记为「已回单」。",
            "「回单导出」按平台导出，便于回传平台。",
        ], "13-首页操作台.png"),
        (7, "对账导出", [
            "按商家或平台统计指定日期范围订单。",
            "一键导出对账 Excel。",
            "「未定义」及未分单商家不出现在列表。",
        ], "02-对账.png"),
        (8, "售后处理", [
            "首页标记售后并填写原因。",
            "售后页支持完结、取消与导出。",
            "默认展示近 30 天数据，可按条件筛选。",
        ], "03-售后.png"),
        (9, "商品价格维护", [
            "维护各平台商品成本价与供货价。",
            "支持表格改价、模板下载与批量导入。",
            "首页据此汇总营业额、成本与利润。",
        ], "05-商品价格维护.png"),
        (10, "回收站", [
            "首页删除的订单进入回收站（软删除）。",
            "支持批量/单条恢复或永久清除。",
            "默认查询近 30 天数据。",
        ], "04-回收站.png"),
        (11, "数据归档", [
            "将历史订单移入归档表，减轻首页数据量。",
            "支持预览与恢复。",
            "与回收站恢复是两套不同机制。",
        ], "10-数据归档.png"),
        (12, "软件授权", [
            "单机版一机一码离线授权。",
            "绑定机器码，通过激活码完成授权。",
            "控制到期与正版使用。",
        ], "11-软件授权.png"),
        (13, "使用手册", [
            "软件内置按模块分节的操作说明。",
            "覆盖首页、对账、售后、配置等场景。",
            "可随时在界面中查阅，降低上手成本。",
        ], "12-使用手册.png"),
    ]

    for step in steps:
        step_slide(prs, *step)

    # 补充字段映射单独页（步骤里合并了，再放一张截图页）
    step_slide(
        prs, 14, "字段映射",
        [
            "自定义界面字段的显示名称。",
            "使列表表头更贴合团队内部叫法。",
            "不影响底层数据存储。",
        ],
        "07-字段映射.png",
    )

    end_slide(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"OK: {OUT_PATH}")
    print(f"size: {OUT_PATH.stat().st_size}")


if __name__ == "__main__":
    build()
